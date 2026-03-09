"""
妆策AI 竞赛答辩PPT生成脚本
9章结构，对标成熟方案（阿里诚聘/跨境电商人才筛选）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

OUT_PATH = r"E:\meiz\zhuangce-ai\zhuangce-ai-最终归档-第一批到第八批\第一批-项目立项层-阶段1\Word版\妆策AI-竞赛答辩PPT-升级版.pptx"

# ── 配色 ──────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 深夜蓝（背景/标题栏）
C_MAIN   = RGBColor(0x16, 0x21, 0x3E)   # 深蓝
C_ACCENT = RGBColor(0x0F, 0x3C, 0x78)   # 中蓝
C_BRIGHT = RGBColor(0x00, 0xB4, 0xD8)   # 亮蓝（强调）
C_GOLD   = RGBColor(0xFF, 0xC3, 0x00)   # 金黄（评分/数据）
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xE8, 0xF4, 0xFD)   # 浅蓝底
C_GRAY   = RGBColor(0x88, 0x88, 0x99)
C_RED    = RGBColor(0xFF, 0x44, 0x44)
C_GREEN  = RGBColor(0x00, 0xC8, 0x80)

W = Inches(13.33)   # 16:9 宽
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 空白布局


# ── 工具函数 ──────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def bg(slide, color=C_DARK):
    """全背景色"""
    add_rect(slide, 0, 0, W, H, fill=color)


def header_bar(slide, title, chapter_num=None, bar_color=C_ACCENT):
    """顶部标题栏"""
    bar_h = Inches(0.9)
    add_rect(slide, 0, 0, W, bar_h, fill=bar_color)
    if chapter_num:
        add_text(slide, str(chapter_num), Inches(0.2), Inches(0.05),
                 Inches(0.6), bar_h, size=28, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(0.85), Inches(0.1),
                 Inches(11), bar_h, size=24, bold=True, color=C_WHITE)
    else:
        add_text(slide, title, Inches(0.35), Inches(0.1),
                 Inches(12), bar_h, size=24, bold=True, color=C_WHITE)


def footer(slide, text="妆策AI · 美妆新零售实战赛"):
    f_h = Inches(0.32)
    add_rect(slide, 0, H - f_h, W, f_h, fill=C_MAIN)
    add_text(slide, text, Inches(0.3), H - f_h,
             Inches(12), f_h, size=10, color=C_GRAY, align=PP_ALIGN.LEFT)


def bullet_list(slide, items, l, t, w, max_h,
                size=15, color=C_WHITE, indent="• ", bold_first=False):
    line_h = Pt(size) * 1.7
    for i, item in enumerate(items):
        prefix = indent if not item.startswith("  ") else "   "
        txt = prefix + item.lstrip()
        add_text(slide, txt, l, t + i * line_h, w, line_h * 1.1,
                 size=size if not item.startswith("  ") else size - 1,
                 color=color,
                 bold=(i == 0 and bold_first))


def kv_pair(slide, key, val, l, t, key_w=Inches(2.2), val_w=Inches(4),
            h=Inches(0.45), key_size=13, val_size=15,
            key_color=C_GRAY, val_color=C_GOLD, bold_val=True):
    add_text(slide, key, l, t, key_w, h, size=key_size, color=key_color)
    add_text(slide, val, l + key_w, t, val_w, h, size=val_size,
             color=val_color, bold=bold_val)


def score_bar(slide, label, score, l, t, max_score=10, bar_max_w=Inches(4)):
    """水平评分条"""
    row_h = Inches(0.38)
    add_text(slide, label, l, t, Inches(2.5), row_h, size=13, color=C_WHITE)
    filled_w = bar_max_w * (score / max_score)
    add_rect(slide, l + Inches(2.6), t + Inches(0.06),
             bar_max_w, row_h - Inches(0.12), fill=C_ACCENT)
    add_rect(slide, l + Inches(2.6), t + Inches(0.06),
             filled_w, row_h - Inches(0.12), fill=C_BRIGHT)
    add_text(slide, f"{score:.1f}", l + Inches(2.6) + filled_w + Inches(0.1),
             t, Inches(0.6), row_h, size=13, color=C_GOLD, bold=True)


def set_notes(slide, text):
    """为幻灯片添加演讲备注（SpeakerNotes）"""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ══════════════════════════════════════════════════
#  SLIDE 1: 封面
# ══════════════════════════════════════════════════
def make_cover(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    # 左侧装饰竖条
    add_rect(slide, 0, 0, Inches(0.18), H, fill=C_BRIGHT)
    # 主标题
    add_text(slide, "妆策AI", Inches(1), Inches(1.2), Inches(11), Inches(1.6),
             size=60, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "面向美妆新零售实战场景的智能推荐与营销转化平台",
             Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             size=22, color=C_BRIGHT, align=PP_ALIGN.LEFT)
    # 分隔线
    add_rect(slide, Inches(1), Inches(3.7), Inches(6), Inches(0.04), fill=C_GOLD)
    # 副信息
    info = [
        "🎯  美妆新零售实战赛 · 校赛参赛作品",
        "🏆  聚焦抖音电商 · DEAR SEED 玫瑰修护洗发水",
        "🤖  15+ AI算法 · 27,512条真实数据 · 34张可视化图表",
    ]
    for i, line in enumerate(info):
        add_text(slide, line, Inches(1), Inches(4.1) + Inches(0.55) * i,
                 Inches(10), Inches(0.5), size=16, color=C_LIGHT)
    # 右下角团队标识
    add_text(slide, "Team Zhuangce AI  ·  2026",
             Inches(9), Inches(6.9), Inches(4), Inches(0.5),
             size=13, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════
#  SLIDE 2: 目录
# ══════════════════════════════════════════════════
def make_toc(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "目录 Table of Contents")
    footer(slide)

    chapters = [
        ("01", "痛点分析", "美妆电商三大核心痛点"),
        ("02", "核心技术", "六维评分 + 五目标优化 + 15+算法"),
        ("03", "平台功能", "爆款预测 / 推荐 / 内容生成 / 仪表盘"),
        ("04", "商业模式", "三方共赢生态 + SaaS授权"),
        ("05", "差异化优势", "vs 竞品对比分析"),
        ("06", "实施计划", "三阶段12个月路线图"),
        ("07", "核心创新", "六维雷达图 + 转化漏斗 + 多目标优化"),
        ("08", "数据可视化", "34张图表 + 21个JSON验证结果"),
        ("09", "总结", "完成度清单 · 核心价值"),
    ]
    col_w = Inches(6.2)
    for i, (num, title, desc) in enumerate(chapters):
        row = i % 5
        col = i // 5
        top = Inches(1.1) + row * Inches(1.18)
        left = Inches(0.4) + col * col_w

        add_rect(slide, left, top, col_w - Inches(0.2), Inches(1.0), fill=C_ACCENT)
        add_rect(slide, left, top, Inches(0.6), Inches(1.0), fill=C_BRIGHT)
        add_text(slide, num, left + Inches(0.05), top + Inches(0.1),
                 Inches(0.5), Inches(0.8), size=20, bold=True,
                 color=C_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, title, left + Inches(0.7), top + Inches(0.05),
                 col_w - Inches(1.1), Inches(0.5), size=16, bold=True, color=C_WHITE)
        add_text(slide, desc, left + Inches(0.7), top + Inches(0.5),
                 col_w - Inches(1.1), Inches(0.45), size=12, color=C_GRAY)


# ══════════════════════════════════════════════════
#  SLIDE 3: 痛点分析
# ══════════════════════════════════════════════════
def make_pain_points(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "痛点分析", "01", C_ACCENT)
    footer(slide)

    pains = [
        ("❶ 选品靠感觉，无数据支撑",
         "运营人员依赖经验选品，缺乏算法预测，选错产品成本高昂",
         "→ 妆策AI 六维评分系统量化传播潜力，GBDT预测AUC=0.9993"),
        ("❷ 内容同质化严重，爆款难复制",
         "洗发水赛道内容千篇一律，缺乏基于数据的差异化内容策略",
         "→ UCB/Thompson最优臂识别，TF-IDF精准提取高权重卖点"),
        ("❸ 运营靠感觉，时机全靠猜",
         "投放时间随机，未能利用平台流量规律，ROI低下",
         "→ ARIMA时序预测，开学季boost=1.15，精准锁定投放窗口"),
    ]

    for i, (title, prob, sol) in enumerate(pains):
        top = Inches(1.1) + i * Inches(1.9)
        # 背景块
        add_rect(slide, Inches(0.4), top, Inches(12.4), Inches(1.75), fill=C_MAIN)
        add_rect(slide, Inches(0.4), top, Inches(0.08), Inches(1.75), fill=C_RED)
        # 内容
        add_text(slide, title, Inches(0.65), top + Inches(0.1),
                 Inches(12), Inches(0.5), size=17, bold=True, color=C_WHITE)
        add_text(slide, prob, Inches(0.65), top + Inches(0.55),
                 Inches(12), Inches(0.45), size=13, color=C_LIGHT)
        add_text(slide, sol, Inches(0.65), top + Inches(1.05),
                 Inches(12), Inches(0.55), size=13, color=C_GREEN, bold=True)


# ══════════════════════════════════════════════════
#  SLIDE 4: 核心技术 - 算法全景
# ══════════════════════════════════════════════════
def make_tech_overview(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "核心技术 · 算法体系全景", "02", C_ACCENT)
    footer(slide)

    modules = [
        ("预测模块", C_BRIGHT, [
            "GBDT 购买预测  AUC=0.9993",
            "ARIMA 时序预测  7天销量曲线",
            "OLS 价格弹性  β=-0.9515",
            "DIN 深度兴趣  购买分离度2.08",
        ]),
        ("推荐模块", RGBColor(0x00, 0xC8, 0x96), [
            "LightGCN 图推荐  HR@10=0.8413",
            "SASRec 序列推荐  HR@10=0.0056",
            "SVD 矩阵分解  解释方差70.3%",
            "ItemCF 协同过滤  命中率28.8%",
        ]),
        ("用户分析", C_GOLD, [
            "K-Means 用户聚类（5类）",
            "RFM 用户分层（4级）",
            "Markov Chain 行为预测",
            "对比学习 轮廓系数+30%",
        ]),
        ("内容优化", RGBColor(0xFF, 0x66, 0x99), [
            "TF-IDF 卖点提取  修护权重13.76",
            "LDA 主题模型  5个内容主题",
            "UCB 多臂老虎机  教程合集最优",
            "Thompson Sampling  策略优化",
        ]),
    ]

    box_w = Inches(3.1)
    for col, (mod_name, mod_color, items) in enumerate(modules):
        left = Inches(0.35) + col * (box_w + Inches(0.12))
        add_rect(slide, left, Inches(1.05), box_w, Inches(0.42), fill=mod_color)
        add_text(slide, mod_name, left + Inches(0.1), Inches(1.08),
                 box_w - Inches(0.2), Inches(0.38),
                 size=15, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        for row, item in enumerate(items):
            top = Inches(1.55) + row * Inches(1.25)
            add_rect(slide, left, top, box_w, Inches(1.15), fill=C_MAIN)
            add_text(slide, item, left + Inches(0.1), top + Inches(0.15),
                     box_w - Inches(0.2), Inches(0.85), size=13, color=C_WHITE, wrap=True)

    add_text(slide, "共 15+ 算法 · 10个分析脚本全部通过 · 34张可视化图表 · 21个JSON验证文件",
             Inches(0.35), Inches(6.8), Inches(12.5), Inches(0.45),
             size=13, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
#  SLIDE 5: 六维评分系统
# ══════════════════════════════════════════════════
def make_six_dim(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "核心技术 · 六维产品传播潜力评分系统", "02", C_ACCENT)
    footer(slide)

    dims = [
        ("D1 内容传播潜力", 7.5, "20%", "UCB/Thompson Sampling", "教程合集转化率3.67%最优"),
        ("D2 购买转化概率", 8.5, "20%", "GBDT购买预测", "AUC=0.9993，精准率极高"),
        ("D3 人群精准匹配", 7.0, "20%", "K-Means + RFM分层", "学生党聚类命中率高"),
        ("D4 时间窗口质量", 8.0, "15%", "ARIMA时序预测", "开学季boost=1.15"),
        ("D5 商品关联强度", 7.5, "15%", "LightGCN + Apriori", "HR@10=0.8413，lift>1.2"),
        ("D6 价格竞争力",   6.5, "10%", "OLS价格弹性回归", "β=-0.9515，100-150元存在压力"),
    ]

    add_text(slide, "DEAR SEED 玫瑰修护洗发水  ·  综合传播潜力评分",
             Inches(0.4), Inches(1.05), Inches(8), Inches(0.45),
             size=14, color=C_GRAY)
    add_text(slide, "7.50 / 10",
             Inches(8.5), Inches(0.88), Inches(4.5), Inches(0.7),
             size=38, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)

    for i, (label, score, weight, algo, evidence) in enumerate(dims):
        top = Inches(1.6) + i * Inches(0.9)
        # 维度名+权重
        add_text(slide, label, Inches(0.4), top + Inches(0.05),
                 Inches(2.4), Inches(0.38), size=13, bold=True, color=C_WHITE)
        add_text(slide, f"权重{weight}", Inches(2.85), top + Inches(0.08),
                 Inches(0.8), Inches(0.3), size=11, color=C_GRAY)
        # 评分条
        bar_total_w = Inches(3.5)
        filled_w    = bar_total_w * (score / 10)
        add_rect(slide, Inches(3.7), top + Inches(0.1), bar_total_w, Inches(0.28), fill=C_ACCENT)
        add_rect(slide, Inches(3.7), top + Inches(0.1), filled_w,    Inches(0.28), fill=C_BRIGHT)
        add_text(slide, f"{score}", Inches(7.3), top + Inches(0.05),
                 Inches(0.5), Inches(0.38), size=14, bold=True, color=C_GOLD)
        # 算法+证据
        add_text(slide, f"{algo}｜{evidence}",
                 Inches(7.9), top + Inches(0.05),
                 Inches(5.0), Inches(0.38), size=11, color=C_GRAY, italic=True)


# ══════════════════════════════════════════════════
#  SLIDE 6: 平台功能
# ══════════════════════════════════════════════════
def make_platform(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "平台功能 · 四大核心页面", "03", C_ACCENT)
    footer(slide)

    pages = [
        ("🔮 爆款预测页", C_BRIGHT, [
            "输入产品名/SKU",
            "六维评分雷达图展示",
            "最佳投放时间建议",
            "风险预警提示",
        ]),
        ("🎯 智能推荐页", RGBColor(0x00, 0xC8, 0x96), [
            "相似商品推荐（LightGCN）",
            "目标人群画像卡片",
            "场景标签匹配",
            "平台策略建议",
        ]),
        ("✍️ 内容生成页", C_GOLD, [
            "抖音短视频标题3条",
            "直播开场/逼单话术",
            "视频hook前3秒脚本",
            "种草文案（小红书风格）",
        ]),
        ("📊 数据仪表盘", RGBColor(0xFF, 0x66, 0x99), [
            "实时运营数据概览",
            "转化漏斗可视化",
            "GMV目标进度追踪",
            "竞赛观测点达成率",
        ]),
    ]

    box_w = Inches(3.0)
    for i, (name, color, items) in enumerate(pages):
        left = Inches(0.35) + i * (box_w + Inches(0.18))
        add_rect(slide, left, Inches(1.0), box_w, Inches(0.5), fill=color)
        add_text(slide, name, left + Inches(0.1), Inches(1.03),
                 box_w - Inches(0.2), Inches(0.44),
                 size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        add_rect(slide, left, Inches(1.5), box_w, Inches(5.5), fill=C_MAIN)
        for j, item in enumerate(items):
            add_text(slide, f"▸  {item}",
                     left + Inches(0.2), Inches(1.65) + j * Inches(0.9),
                     box_w - Inches(0.35), Inches(0.75),
                     size=13, color=C_WHITE, wrap=True)

    # API说明
    add_text(slide, "技术栈：Flask后端 + Vue3 + ECharts + Element Plus  |  5个核心API: /predict  /recommend  /content/generate  /dashboard",
             Inches(0.35), Inches(7.05), Inches(12.5), Inches(0.38),
             size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
#  SLIDE 7: 商业模式 - 三方共赢
# ══════════════════════════════════════════════════
def make_business(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "商业模式 · 三方共赢生态", "04", C_ACCENT)
    footer(slide)

    # 三方卡片
    parties = [
        ("品牌方\n（DEAR SEED等）", C_BRIGHT, [
            "AI六维评分筛选高潜产品",
            "TF-IDF精准提取卖点",
            "ARIMA预测最佳推广时间",
            "→ 降低试错成本，提升ROI",
        ]),
        ("平台\n（抖音电商）", RGBColor(0x00, 0xC8, 0x96), [
            "差异化高质量内容供给",
            "LightGCN图推荐提升转化",
            "K-Means精准用户分群",
            "→ 提升用户停留时长和GMV",
        ]),
        ("消费者\n（学生党女性）", C_GOLD, [
            "SVD个性化推荐匹配真需求",
            "价格竞争力评分找高性价比",
            "AI内容质量高，告别劣质广告",
            "→ 改善购物体验，降低决策成本",
        ]),
    ]
    box_w = Inches(3.9)
    for i, (name, color, items) in enumerate(parties):
        left = Inches(0.35) + i * (box_w + Inches(0.2))
        add_rect(slide, left, Inches(1.05), box_w, Inches(5.7), fill=C_MAIN)
        add_rect(slide, left, Inches(1.05), box_w, Inches(0.9), fill=color)
        add_text(slide, name, left + Inches(0.15), Inches(1.1),
                 box_w - Inches(0.3), Inches(0.8),
                 size=16, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            top = Inches(2.1) + j * Inches(1.1)
            add_text(slide, item, left + Inches(0.2), top,
                     box_w - Inches(0.4), Inches(1.0),
                     size=13, color=C_WHITE if j < 3 else C_GREEN,
                     bold=(j == 3), wrap=True)

    # 收入模型
    add_rect(slide, Inches(0.35), Inches(6.8), Inches(12.5), Inches(0.5), fill=C_ACCENT)
    add_text(slide, "收入模型：SaaS年费（品牌方订阅）+ 竞赛阶段GMV分成（每周≥300元带货目标）",
             Inches(0.5), Inches(6.83), Inches(12.2), Inches(0.44),
             size=13, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
#  SLIDE 8: 差异化优势
# ══════════════════════════════════════════════════
def make_advantage(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "差异化优势 · vs 竞品对比", "05", C_ACCENT)
    footer(slide)

    headers = ["能力维度", "传统电商运营", "通用内容工具", "妆策AI（本项目）"]
    rows = [
        ["选品决策依据", "经验/感觉", "关键词热度", "六维算法评分（AUC=0.9993）"],
        ["内容策略来源", "人工撰写", "模板套用", "UCB/Thompson Sampling最优策略"],
        ["投放时机判断", "运营感知", "平台建议", "ARIMA时序预测（boost=1.15）"],
        ["用户画像来源", "粗糙分层", "无", "K-Means聚类+RFM精细分层"],
        ["推荐逻辑",     "同类热销", "协同过滤", "LightGCN图神经网络（HR@10=0.84）"],
        ["数据基础",     "无",        "部分",      "27,512条真实销售数据+66MB行为数据"],
        ["算法可解释性", "无",        "低",         "每维度有算法溯源，可追溯可复现"],
    ]

    colors_col = [C_GRAY, C_RED, C_GRAY, C_GREEN]
    col_widths  = [Inches(2.5), Inches(3.0), Inches(3.0), Inches(3.8)]
    col_starts  = [Inches(0.35), Inches(2.9), Inches(5.95), Inches(9.0)]
    row_h = Inches(0.73)

    for ci, (hdr, cw, cl_start, col_color) in enumerate(zip(headers, col_widths, col_starts, colors_col)):
        add_rect(slide, cl_start, Inches(1.0), cw - Inches(0.06), Inches(0.5),
                 fill=C_ACCENT if ci == 3 else C_MAIN)
        add_text(slide, hdr, cl_start + Inches(0.08), Inches(1.02),
                 cw - Inches(0.15), Inches(0.46),
                 size=13, bold=True, color=C_GOLD if ci == 3 else C_WHITE,
                 align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        top = Inches(1.55) + ri * row_h
        for ci, (cell, cw, cl_start) in enumerate(zip(row, col_widths, col_starts)):
            fill_color = C_ACCENT if ci == 3 else (C_MAIN if ri % 2 == 0 else RGBColor(0x14, 0x1E, 0x36))
            add_rect(slide, cl_start, top, cw - Inches(0.06), row_h - Inches(0.05), fill=fill_color)
            add_text(slide, cell, cl_start + Inches(0.1), top + Inches(0.05),
                     cw - Inches(0.2), row_h - Inches(0.1),
                     size=12, color=C_GOLD if ci == 3 else C_WHITE, wrap=True)


# ══════════════════════════════════════════════════
#  SLIDE 9: 实施计划
# ══════════════════════════════════════════════════
def make_roadmap(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "实施计划 · 三阶段路线图", "06", C_ACCENT)
    footer(slide)

    phases = [
        ("Phase 1\n技术验证", C_BRIGHT, "已完成 ✓", [
            "10个数据分析脚本全部通过",
            "15+算法实现并验证",
            "27,512条真实数据处理完成",
            "34张图表+21个JSON输出",
            "MVP产品确认（DEAR SEED）",
            "竞赛评分适配方案制定",
        ]),
        ("Phase 2\n产品开发", RGBColor(0x00, 0xC8, 0x96), "进行中（5-8天）", [
            "Flask后端5个核心API开发",
            "Vue3+ECharts前端页面完善",
            "前后端联调部署",
            "六维评分雷达图可视化",
            "内容生成模块集成",
            "答辩PPT制作",
        ]),
        ("Phase 3\n实战运营", C_GOLD, "待启动（竞赛期间）", [
            "15个账号开通并绑定赛事平台",
            "签订合作协议（必须项）",
            "每周5条短视频/图文发布",
            "每周1场有效直播（≥2小时）",
            "每周GMV目标≥300元",
            "竞赛观测点持续达成",
        ]),
    ]

    box_w = Inches(4.0)
    for i, (name, color, status, items) in enumerate(phases):
        left = Inches(0.35) + i * (box_w + Inches(0.2))
        add_rect(slide, left, Inches(1.0), box_w, Inches(0.6), fill=color)
        add_text(slide, name, left + Inches(0.1), Inches(1.02),
                 box_w * 0.6, Inches(0.55),
                 size=16, bold=True, color=C_DARK)
        add_text(slide, status, left + box_w * 0.6, Inches(1.1),
                 box_w * 0.38, Inches(0.4),
                 size=11, color=C_DARK, align=PP_ALIGN.RIGHT)
        add_rect(slide, left, Inches(1.6), box_w, Inches(5.15), fill=C_MAIN)
        for j, item in enumerate(items):
            add_text(slide, f"  ▸  {item}",
                     left + Inches(0.1), Inches(1.7) + j * Inches(0.82),
                     box_w - Inches(0.2), Inches(0.75),
                     size=13, color=C_WHITE, wrap=True)


# ══════════════════════════════════════════════════
#  SLIDE 10: 核心创新
# ══════════════════════════════════════════════════
def make_innovation(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "核心创新 · 五大创新点", "07", C_ACCENT)
    footer(slide)

    innovations = [
        ("①", "六维产品传播潜力评分",
         "将选品决策从「感觉判断」升级为「算法量化」，每维度均有真实算法证据支撑，可解释、可追溯、可复现。DEAR SEED综合评分7.50/10。",
         C_BRIGHT),
        ("②", "五目标内容投放多目标优化",
         "同时优化关键词匹配(35%)+时间窗口(25%)+制作成本(20%)+购买意愿(10%)+价格匹配(10%)，每周生成最优内容日历。",
         RGBColor(0x00, 0xC8, 0x96)),
        ("③", "Markov Chain用户转化漏斗",
         "基于用户行为序列的马尔可夫链建模，识别浏览→加购→购买的最优引导路径，精准干预不同阶段用户。",
         C_GOLD),
        ("④", "AI端到端内容生成闭环",
         "TF-IDF提取卖点→UCB选择最优内容形式→AI生成抖音标题/直播话术/视频Hook，实现「数据驱动内容」完整闭环。",
         RGBColor(0xFF, 0x66, 0x99)),
        ("⑤", "图神经网络商品关联推荐",
         "LightGCN构建商品图谱，HR@10=0.8413，NDCG@10=0.7031，精准识别修护洗发与护发精华的跨品类关联带货机会。",
         RGBColor(0xAA, 0x44, 0xFF)),
    ]

    for i, (num, title, desc, color) in enumerate(innovations):
        row = i % 3
        col = i // 3
        left = Inches(0.35) + col * Inches(6.6)
        top  = Inches(1.05) + row * Inches(2.05)
        w    = Inches(6.3)
        h    = Inches(1.9)
        add_rect(slide, left, top, w, h, fill=C_MAIN)
        add_rect(slide, left, top, Inches(0.55), h, fill=color)
        add_text(slide, num, left + Inches(0.05), top + Inches(0.6),
                 Inches(0.45), Inches(0.65),
                 size=22, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, title, left + Inches(0.65), top + Inches(0.1),
                 w - Inches(0.75), Inches(0.48),
                 size=15, bold=True, color=C_WHITE)
        add_text(slide, desc, left + Inches(0.65), top + Inches(0.58),
                 w - Inches(0.75), Inches(1.25),
                 size=12, color=C_LIGHT, wrap=True)


# ══════════════════════════════════════════════════
#  SLIDE 11: 数据可视化 - 关键指标展示
# ══════════════════════════════════════════════════
def make_data_viz(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "数据可视化 · 关键算法指标", "08", C_ACCENT)
    footer(slide)

    metrics = [
        ("GBDT购买预测", "AUC", "0.9993", C_BRIGHT),
        ("LightGCN图推荐", "HR@10", "0.8413", RGBColor(0x00, 0xC8, 0x96)),
        ("LightGCN图推荐", "NDCG@10", "0.7031", RGBColor(0x00, 0xC8, 0x96)),
        ("UCB最优内容", "最优转化率", "3.67%", C_GOLD),
        ("对比学习增强", "轮廓系数↑", "+30.01%", RGBColor(0xFF, 0x66, 0x99)),
        ("ARIMA时序", "开学季boost", "×1.15", C_BRIGHT),
        ("OLS价格弹性", "弹性系数", "-0.9515", C_GOLD),
        ("SVD矩阵分解", "解释方差", "70.27%", RGBColor(0xAA, 0x44, 0xFF)),
    ]

    card_w = Inches(3.0)
    card_h = Inches(1.5)
    for i, (algo, metric, value, color) in enumerate(metrics):
        col = i % 4
        row = i // 4
        left = Inches(0.35) + col * (card_w + Inches(0.18))
        top  = Inches(1.05) + row * (card_h + Inches(0.25))
        add_rect(slide, left, top, card_w, card_h, fill=C_MAIN)
        add_rect(slide, left, top, card_w, Inches(0.08), fill=color)
        add_text(slide, algo, left + Inches(0.12), top + Inches(0.15),
                 card_w - Inches(0.2), Inches(0.4),
                 size=12, color=C_GRAY)
        add_text(slide, value, left + Inches(0.12), top + Inches(0.5),
                 card_w - Inches(0.2), Inches(0.65),
                 size=30, bold=True, color=C_GOLD)
        add_text(slide, metric, left + Inches(0.12), top + Inches(1.15),
                 card_w - Inches(0.2), Inches(0.3),
                 size=12, color=C_LIGHT)

    add_text(slide, "完整验证：34张可视化图表  +  21个JSON输出文件  +  27,512条真实销售数据",
             Inches(0.35), Inches(6.9), Inches(12.5), Inches(0.45),
             size=13, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
#  SLIDE 12: 总结 - 完成度清单
# ══════════════════════════════════════════════════
def make_summary(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    header_bar(slide, "总结 · 核心价值与完成度", "09", C_ACCENT)
    footer(slide)

    # 左侧：完成度清单
    add_rect(slide, Inches(0.35), Inches(1.0), Inches(6.5), Inches(5.9), fill=C_MAIN)
    add_text(slide, "✅  完成度清单",
             Inches(0.5), Inches(1.05), Inches(6), Inches(0.5),
             size=16, bold=True, color=C_BRIGHT)

    items_done = [
        "✓  真实场景分析（美妆新零售·抖音电商）",
        "✓  用户画像确认（学生党·18-22岁·100-150元）",
        "✓  痛点挖掘（同质化内容/无算法选品/缺时序策略）",
        "✓  算法设计与实现（15+算法·10个脚本·700+行代码）",
        "✓  真实数据验证（27,512条销售+66MB行为数据）",
        "✓  运行结果验证（34张图表+21个JSON输出文件）",
        "✓  六维评分框架（每维度有算法证据，可追溯）",
        "✓  内容运营方案（每周5条短视频+1场直播+300GMV）",
        "✓  竞赛评分适配（创新+创意强覆盖，运营待启动）",
    ]
    for i, item in enumerate(items_done):
        add_text(slide, item,
                 Inches(0.5), Inches(1.65) + i * Inches(0.57),
                 Inches(6.2), Inches(0.52),
                 size=12, color=C_GREEN if item.startswith("✓") else C_LIGHT)

    # 右侧：核心价值
    add_rect(slide, Inches(7.1), Inches(1.0), Inches(5.85), Inches(5.9), fill=C_MAIN)
    add_text(slide, "💡  核心价值",
             Inches(7.25), Inches(1.05), Inches(5.5), Inches(0.5),
             size=16, bold=True, color=C_GOLD)

    values = [
        ("解决真实问题",
         "运营靠感觉 → AI六维评分\n内容同质化 → 最优策略推荐\n时机全靠猜 → ARIMA时序预测"),
        ("技术创新",
         "六维产品评分 · 五目标优化\nMarkov转化漏斗 · 图神经推荐\nAI端到端内容生成闭环"),
        ("可立即实施",
         "完整代码已验证运行\n商业模式清晰可落地\n三阶段实施路线已规划"),
    ]
    for i, (title, desc) in enumerate(values):
        top = Inches(1.65) + i * Inches(1.7)
        add_rect(slide, Inches(7.15), top, Inches(5.7), Inches(1.55),
                 fill=RGBColor(0x14, 0x1E, 0x36))
        add_rect(slide, Inches(7.15), top, Inches(0.06), Inches(1.55), fill=C_GOLD)
        add_text(slide, title, Inches(7.3), top + Inches(0.1),
                 Inches(5.4), Inches(0.42),
                 size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, Inches(7.3), top + Inches(0.5),
                 Inches(5.4), Inches(1.0),
                 size=12, color=C_LIGHT, wrap=True)

    # 底部口号
    add_rect(slide, Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.38), fill=C_ACCENT)
    add_text(slide, "「用数据说话，让AI为美妆新零售赋能」",
             Inches(0.5), Inches(7.03), Inches(12.2), Inches(0.34),
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
#  SLIDE 13: Q&A
# ══════════════════════════════════════════════════
def make_qa(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, C_DARK)
    add_rect(slide, 0, 0, Inches(0.18), H, fill=C_BRIGHT)
    add_text(slide, "Q & A", Inches(1), Inches(2.2), Inches(11), Inches(2.0),
             size=90, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Thank you for listening!",
             Inches(1), Inches(4.3), Inches(11), Inches(0.7),
             size=24, color=C_BRIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "妆策AI · 美妆新零售实战赛 · 校赛参赛作品",
             Inches(1), Inches(5.3), Inches(11), Inches(0.5),
             size=15, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
#  BUILD ALL SLIDES
# ══════════════════════════════════════════════════
make_cover(prs)
set_notes(prs.slides[-1], "各位老师好，我们团队参赛项目是「妆策AI」——面向美妆新零售实战场景的智能推荐与营销转化平台。项目聚焦抖音电商赛道，以DEAR SEED玫瑰修护洗发水为MVP示例产品，已完成第一阶段技术验证：15+AI算法、27,512条真实数据、34张可视化图表、21个JSON验证文件。请各位老师评委多多指导。")

make_toc(prs)
set_notes(prs.slides[-1], "本次答辩共分9个部分，逻辑链为：市场痛点→技术创新→产品落地→商业价值→数据验证。整个演讲约12-15分钟，之后留5分钟Q&A环节。")

make_pain_points(prs)
set_notes(prs.slides[-1], "美妆电商运营面临三大核心痛点：第一，选品完全靠经验，试错成本极高——妆策AI用六维算法评分量化传播潜力，GBDT预测AUC=0.9993；第二，内容同质化严重，爆款无法系统复制——UCB/Thompson最优臂识别最优内容形式，教程合集转化率3.67%最高；第三，投放时机随意，ROI低下——ARIMA时序预测，开学季boost=1.15，精准锁定高流量窗口。")

make_tech_overview(prs)
set_notes(prs.slides[-1], "我们构建了由15+算法组成的四大模块，每个算法都已实际运行通过并产出验证结果。预测模块的GBDT AUC=0.9993，接近完美；推荐模块的LightGCN HR@10=0.8413，性能优秀；用户分析的对比学习让用户表示质量提升30%；内容优化的UCB识别出教程合集为最优内容形式。10个脚本全部通过，共生成34张图表+21个JSON验证文件。")

make_six_dim(prs)
set_notes(prs.slides[-1], "六维评分系统是本项目的核心创新。以DEAR SEED为例，综合评分7.50/10。购买转化概率最高8.5分（GBDT AUC=0.9993，精准率极高）；时间窗口8.0分（开学季是最佳投放期）；内容传播和商品关联均7.5分；人群精准匹配7.0分；价格竞争力最低6.5分，说明100-150元价位在学生党群体中存在一定溢价压力，需要内容价值支撑。每个分数都有算法证据，可在代码中追溯。")

make_platform(prs)
set_notes(prs.slides[-1], "平台共4个核心页面：爆款预测页——输入产品名即可获得六维雷达图评分和最佳投放时间；智能推荐页——LightGCN图推荐提供相似商品和人群画像；内容生成页——一键输出3条抖音标题、直播话术和视频Hook前3秒脚本；数据仪表盘——实时追踪运营数据、GMV目标和竞赛观测点。技术栈：Flask后端+Vue3+ECharts，5个标准API接口，第二阶段将完成开发。")

make_business(prs)
set_notes(prs.slides[-1], "商业模式是三方共赢的生态体系。品牌方获得AI选品和内容策略降低试错成本；抖音电商获得高质量内容提升转化率；学生党消费者获得精准个性化推荐改善购物体验。收入模型分两条线：产品正式落地后收SaaS年费；竞赛阶段以每周GMV≥300元为运营基准，每周5条短视频+1场直播，直接对应竞赛创业业绩评分。")

make_advantage(prs)
set_notes(prs.slides[-1], "这张对比表直接展示妆策AI的差异化优势。横向对比：传统运营几乎全靠人工；通用内容工具有一定帮助但缺乏算法深度；妆策AI每个维度都有算法支撑。最关键的差异在于算法可解释性——我们每一个数字都能追溯到具体脚本和JSON输出文件，这是竞赛答辩中最有说服力的证据。")

make_roadmap(prs)
set_notes(prs.slides[-1], "三阶段实施计划：Phase1技术验证已全部完成，10个脚本+15+算法+34张图表+21个JSON，这是今天答辩的核心内容；Phase2产品开发正在进行，Flask+Vue3全栈开发，预计5-8天完成；Phase3实战运营将在比赛期间启动，15个账号注册绑定赛事平台是最高优先级，每周5条内容+1场直播+GMV≥300元直接支撑创业过程20分和创业业绩15分评分。")

make_innovation(prs)
set_notes(prs.slides[-1], "五大核心创新点总结：第一，六维评分系统——选品从主观判断变为客观量化，每维度有算法证据；第二，五目标内容优化——同时平衡五个相互制约的目标，生成最优内容日历；第三，Markov Chain转化漏斗——识别用户所在行为阶段（浏览/收藏/加购/购买），精准干预；第四，AI端到端内容生成——从选品卖点到内容脚本全流程AI赋能；第五，LightGCN图推荐——HR@10=0.8413，精准发现跨品类关联带货机会。")

make_data_viz(prs)
set_notes(prs.slides[-1], "这8个指标是算法验证的核心数据，每一个都来自真实运行的代码输出，可提供对应的JSON文件和Python脚本供评委查阅。重点说明：GBDT AUC=0.9993接近完美预测；LightGCN HR@10=0.8413在推荐领域属于优秀水平；UCB教程合集转化率3.67%为所有内容形式中最高；对比学习用户表示质量提升30.01%；完整验证材料：34张图表+21个JSON文件，全部可提供。")

make_summary(prs)
set_notes(prs.slides[-1], "总结：第一阶段9项核心工作全部完成，技术验证扎实充分。核心价值三点：解决真实市场痛点（三个核心问题有对应算法方案）；技术创新（五大创新点，15+算法，均有数据验证）；可落地实施（完整代码+清晰商业模式+三阶段路线图）。下一步最紧迫：Phase3账号注册和实操启动，直接支撑竞赛评分中最空缺的30分（创业过程20分+创业业绩15分中的运营部分）。感谢各位老师！")

make_qa(prs)
set_notes(prs.slides[-1], "Q&A备用问答要点：1）算法为什么选GBDT？——AUC=0.9993，对购买预测的非线性特征处理能力强，适合电商场景；2）LightGCN的训练数据来源？——来自66MB美妆用户行为数据集，包含真实用户-商品交互关系；3）六维评分如何融合？——加权均值，权重设计参考各维度对最终转化率的贡献度；4）GMV300元如何实现？——抖音橱窗+直播间+三创好物多渠道并行，AI内容策略提升转化效率；5）创业准备分如何保证？——合作协议签订+15账号注册绑定为最高优先级行动项，已列入Phase3第一周任务。")

prs.save(OUT_PATH)
print(f"PPT saved: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(OUT_PATH):,} bytes")
